from pathlib import Path
import sys

sys.path.insert(0, "/private/tmp/chapel-pptx")

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
OUTPUT = ROOT / "docs" / "presentations" / "chapel-scheduler-august-8-2026.pptx"

WINE = RGBColor(111, 23, 37)
WINE_DARK = RGBColor(72, 16, 26)
WINE_PALE = RGBColor(245, 233, 235)
GOLD = RGBColor(184, 138, 53)
GOLD_PALE = RGBColor(248, 238, 206)
INK = RGBColor(35, 32, 30)
MUTED = RGBColor(103, 96, 91)
CREAM = RGBColor(246, 241, 233)
PAPER = RGBColor(255, 253, 249)
LINE = RGBColor(221, 213, 205)
GREEN = RGBColor(57, 114, 88)
GREEN_PALE = RGBColor(234, 245, 239)
WHITE = RGBColor(255, 255, 255)


prs = Presentation()
prs.slide_width = Inches(13.333333)
prs.slide_height = Inches(7.5)
blank = prs.slide_layouts[6]


def rect(slide, x, y, w, h, fill, line=None, radius=True):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    s = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    s.line.color.rgb = line if line else fill
    if radius:
        try:
            s.adjustments[0] = 0.08
        except Exception:
            pass
    return s


def line(slide, x1, y1, x2, y2, color=LINE, width=1.5):
    c = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2)
    )
    c.line.color.rgb = color
    c.line.width = Pt(width)
    return c


def text_box(slide, x, y, w, h, text, size=18, color=INK, bold=False,
             font="Aptos", align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.TOP,
             margin=0.04, italic=False):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    box.text_frame.clear()
    box.text_frame.margin_left = Inches(margin)
    box.text_frame.margin_right = Inches(margin)
    box.text_frame.margin_top = Inches(margin)
    box.text_frame.margin_bottom = Inches(margin)
    box.text_frame.vertical_anchor = valign
    p = box.text_frame.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.name = font
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    return box


def rich_text(slide, x, y, w, h, paragraphs, margin=0.05, spacing=6):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.margin_left = Inches(margin)
    tf.margin_right = Inches(margin)
    tf.margin_top = Inches(margin)
    tf.margin_bottom = Inches(margin)
    for i, item in enumerate(paragraphs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item.get("text", "")
        p.font.name = item.get("font", "Aptos")
        p.font.size = Pt(item.get("size", 17))
        p.font.bold = item.get("bold", False)
        p.font.color.rgb = item.get("color", INK)
        p.level = item.get("level", 0)
        p.space_after = Pt(item.get("after", spacing))
        p.alignment = item.get("align", PP_ALIGN.LEFT)
        if item.get("bullet", False):
            p.text = "•  " + p.text
    return box


def circle_badge(slide, x, y, d, label, fill=WINE, size=18):
    s = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(d), Inches(d))
    s.fill.solid(); s.fill.fore_color.rgb = fill
    s.line.color.rgb = fill
    text_box(slide, x, y, d, d, label, size=size, color=WHITE, bold=True,
             align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE, margin=0)
    return s


def top_rule(slide):
    rect(slide, 0, 0, 13.333, 0.10, GOLD, GOLD, radius=False)


def footer(slide, number):
    text_box(slide, 0.55, 7.11, 5.7, 0.2, "CHAPEL SCHEDULER  ·  AUGUST 8, 2026",
             size=8, color=MUTED, bold=True)
    text_box(slide, 12.15, 7.08, 0.6, 0.24, str(number), size=9, color=MUTED,
             bold=True, align=PP_ALIGN.RIGHT)


def header(slide, title, subtitle=None, number=None):
    top_rule(slide)
    text_box(slide, 0.58, 0.34, 12.0, 0.58, title, size=29, color=WINE_DARK,
             bold=False, font="Georgia")
    if subtitle:
        text_box(slide, 0.61, 0.94, 11.6, 0.38, subtitle, size=13, color=MUTED)
    if number is not None:
        footer(slide, number)


def card(slide, x, y, w, h, title, body, badge=None, fill=PAPER, accent=WINE,
         body_size=14):
    rect(slide, x, y, w, h, fill, LINE)
    if badge:
        circle_badge(slide, x + 0.22, y + 0.22, 0.46, badge, accent, 13)
        tx = x + 0.82
        tw = w - 1.04
    else:
        rect(slide, x, y, 0.07, h, accent, accent, radius=False)
        tx = x + 0.28
        tw = w - 0.5
    text_box(slide, tx, y + 0.20, tw, 0.34, title, size=16, color=accent, bold=True)
    text_box(slide, tx, y + 0.67, tw, h - 0.85, body, size=body_size, color=INK)


def notes(slide, body):
    try:
        tf = slide.notes_slide.notes_text_frame
        tf.text = body
    except Exception:
        pass


# 1 — Title
slide = prs.slides.add_slide(blank)
rect(slide, 0, 0, 13.333, 7.5, CREAM, CREAM, radius=False)
rect(slide, 0, 0, 4.05, 7.5, WINE_DARK, WINE_DARK, radius=False)
rect(slide, 4.05, 0, 0.10, 7.5, GOLD, GOLD, radius=False)
circle_badge(slide, 1.34, 1.02, 1.33, "✠", WINE, 34)
text_box(slide, 0.7, 2.67, 2.65, 0.42, "OUR LADY OF VICTORY", 11, GOLD_PALE,
         True, align=PP_ALIGN.CENTER)
text_box(slide, 4.78, 1.52, 7.6, 0.8, "Chapel Scheduler", 38, WINE_DARK,
         False, "Georgia")
text_box(slide, 4.82, 2.48, 7.1, 0.72,
         "One coordinated calendar for services, ministries, and communication",
         22, INK, False, "Georgia")
line(slide, 4.84, 3.55, 11.8, 3.55, GOLD, 2.2)
text_box(slide, 4.82, 3.88, 6.8, 0.42, "PLANNING MEETING  ·  AUGUST 8, 2026",
         12, WINE, True)
text_box(slide, 4.82, 4.53, 6.8, 0.62,
         "Proposal and proof-of-concept review", 17, MUTED)
rect(slide, 4.82, 5.52, 6.9, 0.78, PAPER, LINE)
text_box(slide, 5.08, 5.72, 6.36, 0.38,
         "Goal: simplify chapel administration without sacrificing privacy.",
         15, GREEN, True)
notes(slide, "Open by explaining that this is a proposal and a proof of concept, not a finished production system or a request to replace everything at once.")


# 2 — Why now
slide = prs.slides.add_slide(blank)
header(slide, "Why change now?", "The chapel is maintaining the same reality in several disconnected places.", 2)
card(slide, 0.62, 1.55, 3.65, 2.0, "Fragmented tools",
     "Google calendars, Ministry Scheduler Pro, Telegram, spreadsheets, websites, and manual public listings.", "1", WINE_PALE, WINE)
card(slide, 4.84, 1.55, 3.65, 2.0, "Repeated work",
     "The same service or change must be entered more than once—and each copy can drift out of date.", "2", GOLD_PALE, GOLD)
card(slide, 9.06, 1.55, 3.65, 2.0, "Low adoption",
     "Volunteers already use messaging apps. A separate, cumbersome scheduling app is easy to ignore.", "3", GREEN_PALE, GREEN)
rect(slide, 1.42, 4.25, 10.5, 1.48, PAPER, LINE)
text_box(slide, 1.78, 4.58, 9.78, 0.35, "The practical result", 15, WINE, True, align=PP_ALIGN.CENTER)
text_box(slide, 1.82, 5.02, 9.7, 0.45,
         "More follow-up, more chances for conflict, and less confidence that any one calendar is current.",
         20, INK, False, "Georgia", PP_ALIGN.CENTER)
notes(slide, "Use concrete examples from current practice, but avoid criticizing the people maintaining these systems. The problem is fragmentation, not effort.")


# 3 — Vision
slide = prs.slides.add_slide(blank)
header(slide, "The proposal: one coordinated source of truth", "Enter a service once; show each person only what they are authorized to see.", 3)
circle_badge(slide, 5.62, 2.18, 2.1, "ONE\nCALENDAR", WINE, 18)
items = [
    (0.65, 1.55, "Public", "Masses, Confessions, Rosary, special services"),
    (0.65, 4.45, "Volunteers", "Assignments, openings, preferences, substitutes"),
    (9.20, 1.55, "Leaders", "Drafts, shortages, qualifications, confirmations"),
    (9.20, 4.45, "Priests & delegates", "Celebrants, private availability, protected appointments"),
]
for x, y, title, body in items:
    card(slide, x, y, 3.45, 1.45, title, body, fill=PAPER,
         accent=GREEN if title == "Public" else WINE, body_size=13)
    line(slide, x + (3.45 if x < 5 else 0), y + 0.72,
         5.62 if x < 5 else 7.72, 3.23, GOLD, 1.7)
text_box(slide, 4.58, 5.35, 4.2, 0.46, "Permissioned views—not separate calendars",
         14, GOLD, True, align=PP_ALIGN.CENTER)
notes(slide, "The master calendar does not mean everyone sees everything. It means one event can safely produce different views for the public, volunteers, leaders, and priests.")


# 4 — Experience by role
slide = prs.slides.add_slide(blank)
header(slide, "Simple for each person", "The system becomes more useful after sign-in without becoming more revealing.", 4)
roles = [
    ("PUBLIC VISITOR", "No account", "Mass times and public service information", GOLD),
    ("ACCOUNT HOLDER", "Optional", "Chosen alerts and visible ways to help", GREEN),
    ("VOLUNTEER", "Approved", "Relevant assignments, openings, and instructions", WINE),
    ("LEADER / ADMIN", "Authorized", "Scheduling, exceptions, and oversight", WINE_DARK),
]
for i, (role, access, body, accent) in enumerate(roles):
    y = 1.55 + i * 1.25
    rect(slide, 0.82, y, 11.65, 0.94, PAPER, LINE)
    rect(slide, 0.82, y, 0.10, 0.94, accent, accent, radius=False)
    text_box(slide, 1.16, y + 0.20, 2.2, 0.28, role, 12, accent, True)
    text_box(slide, 3.44, y + 0.20, 1.4, 0.28, access, 12, MUTED, True)
    text_box(slide, 5.05, y + 0.15, 6.95, 0.40, body, 16, INK)
text_box(slide, 1.15, 6.58, 10.8, 0.32,
         "Private calendars and APR details are invisible—not merely shown as locked tabs.",
         13, WINE, True, align=PP_ALIGN.CENTER)
notes(slide, "Emphasize invisible permissions: unauthorized users should not even learn that Father has a restricted appointment calendar or that another volunteer has an APR issue.")


# 5 — Engagement funnel
slide = prs.slides.add_slide(blank)
header(slide, "Make participation easy—but always optional", "The public calendar can become the front door to communication and service.", 5)
steps = [
    ("1", "Check the calendar", "No account required"),
    ("2", "Choose updates", "Every category opt-in"),
    ("3", "See ways to help", "No commitment"),
    ("4", "Leader follows up", "One-time or ongoing"),
]
for i, (num, title, body) in enumerate(steps):
    x = 0.62 + i * 3.18
    circle_badge(slide, x + 0.93, 1.53, 0.72, num, WINE if i in (0, 3) else GOLD, 18)
    if i < 3:
        line(slide, x + 1.65, 1.89, x + 3.0, 1.89, GOLD, 2.0)
    text_box(slide, x, 2.48, 2.56, 0.34, title, 16, WINE_DARK, True,
             align=PP_ALIGN.CENTER)
    text_box(slide, x, 2.91, 2.56, 0.35, body, 12, MUTED,
             align=PP_ALIGN.CENTER)
rect(slide, 1.12, 4.03, 11.08, 1.52, GREEN_PALE, GREEN)
text_box(slide, 1.52, 4.32, 10.3, 0.32, "Examples people would see", 13, GREEN, True,
         align=PP_ALIGN.CENTER)
text_box(slide, 1.52, 4.82, 10.3, 0.38,
         "Liturgical ministries  ·  League cleaning and meals  ·  Maintenance  ·  Events",
         17, INK, False, "Georgia", PP_ALIGN.CENTER)
text_box(slide, 2.12, 5.92, 9.1, 0.40,
         "A name is requested only when someone offers to help.",
         15, WINE, True, align=PP_ALIGN.CENTER)
notes(slide, "Point out the public event link: it says only that opportunities exist. It does not reveal how many people are missing or which ministry needs help.")


# 6 — Scheduling workflow
slide = prs.slides.add_slide(blank)
header(slide, "Less chasing; more predictable scheduling", "Automation proposes. Authorized leaders remain in control.", 6)
labels = [
    ("Preferences", "When people usually can serve"),
    ("Balanced draft", "Qualifications, workload, and conflicts"),
    ("Leader review", "Shortages and overrides stay visible"),
    ("Confirm & remind", "Email first; Telegram can follow"),
    ("Substitute", "Qualified people, private requests"),
]
for i, (title, body) in enumerate(labels):
    x = 0.4 + i * 2.58
    rect(slide, x, 2.0, 2.23, 2.08, PAPER, LINE)
    circle_badge(slide, x + 0.78, 1.58, 0.66, str(i + 1), WINE if i % 2 == 0 else GOLD, 16)
    text_box(slide, x + 0.17, 2.35, 1.89, 0.38, title, 15, WINE_DARK, True,
             align=PP_ALIGN.CENTER)
    text_box(slide, x + 0.22, 2.95, 1.79, 0.72, body, 12, MUTED,
             align=PP_ALIGN.CENTER)
    if i < 4:
        line(slide, x + 2.23, 3.04, x + 2.55, 3.04, GOLD, 2.0)
rect(slide, 1.22, 4.78, 10.88, 1.05, WINE_PALE, WINE)
text_box(slide, 1.62, 5.06, 10.08, 0.45,
         "The leader can always correct the draft; privacy and APR eligibility cannot be overridden.",
         16, WINE, True, align=PP_ALIGN.CENTER)
notes(slide, "This is the operational heart of the proposal. Balance workload and use preferences, but preserve human judgment. Hard privacy and APR rules remain non-overridable.")


# 7 — Privacy
slide = prs.slides.add_slide(blank)
header(slide, "Privacy first—by design", "Collect less, separate permissions, and avoid exposing sensitive status.", 7)
privacy = [
    ("Public", "Services and public liturgical details only", "No volunteer names or staffing counts"),
    ("Minors", "Birth year rather than full birth date", "Guardian visibility and controls"),
    ("APR", "Eligibility metadata only", "No reports, reasons, or signed forms by default"),
    ("Father", "Private details for approved delegates", "Others see only unavailable"),
]
for i, (title, body, detail) in enumerate(privacy):
    x = 0.7 + (i % 2) * 6.15
    y = 1.52 + (i // 2) * 2.42
    rect(slide, x, y, 5.77, 1.88, PAPER, LINE)
    circle_badge(slide, x + 0.25, y + 0.27, 0.58, "✓", GREEN, 16)
    text_box(slide, x + 1.04, y + 0.23, 4.25, 0.34, title, 17, WINE_DARK, True)
    text_box(slide, x + 1.04, y + 0.75, 4.3, 0.34, body, 14, INK)
    text_box(slide, x + 1.04, y + 1.22, 4.3, 0.28, detail, 11, MUTED)
text_box(slide, 1.42, 6.42, 10.5, 0.33,
         "Convenience never grants broader access than a person's actual responsibility.",
         14, WINE, True, align=PP_ALIGN.CENTER)
notes(slide, "Avoid describing implementation details. The message is that privacy is structural: public, ministry, APR, and private priest information are distinct.")


# 8 — Alpha
slide = prs.slides.add_slide(blank)
header(slide, "Start small: the sacristan alpha", "A private, reversible test before normal chapel use.", 8)
rect(slide, 0.72, 1.48, 5.92, 4.72, GREEN_PALE, GREEN)
text_box(slide, 1.10, 1.81, 5.18, 0.38, "IN THE ALPHA", 13, GREEN, True)
rich_text(slide, 1.08, 2.42, 5.15, 3.28, [
    {"text": "Public service calendar", "bullet": True, "size": 16},
    {"text": "Secure sacristan accounts", "bullet": True, "size": 16},
    {"text": "Preferences, absences, and balanced drafts", "bullet": True, "size": 16},
    {"text": "Confirmations, reminders, and substitutes", "bullet": True, "size": 16},
    {"text": "Administrator editing, printing, and history", "bullet": True, "size": 16},
    {"text": "Email notifications and tested backups", "bullet": True, "size": 16},
])
rect(slide, 6.96, 1.48, 5.64, 4.72, PAPER, LINE)
text_box(slide, 7.34, 1.81, 4.88, 0.38, "NOT REQUIRED FOR ALPHA", 13, WINE, True)
rich_text(slide, 7.32, 2.42, 4.86, 3.28, [
    {"text": "Replacing every existing calendar", "bullet": True, "size": 16},
    {"text": "Altar servers, ushers, or other ministries", "bullet": True, "size": 16},
    {"text": "Advanced APR document workflows", "bullet": True, "size": 16},
    {"text": "Website and public-feed integration", "bullet": True, "size": 16},
    {"text": "SMS or full Telegram automation", "bullet": True, "size": 16},
    {"text": "Multiple chapels", "bullet": True, "size": 16},
])
notes(slide, "The point of alpha is to test whether eight sacristans and the administrator can use the core reliably. Everything else is deliberately deferred.")


# 9 — Roadmap
slide = prs.slides.add_slide(blank)
header(slide, "Phased rollout", "Each stage earns the next one through actual use and feedback.", 9)
line(slide, 1.25, 3.15, 12.1, 3.15, LINE, 4)
milestones = [
    (1.0, "NOW", "Design + POC", WINE),
    (3.62, "AUG 8", "Stakeholder decision", GOLD),
    (6.25, "EARLY SEP", "Private sacristan alpha", GREEN),
    (8.88, "LATE OCT", "Servers + ushers pilots", WINE),
    (11.1, "ADVENT", "Normal use if validated", WINE_DARK),
]
for x, date, label, color in milestones:
    circle_badge(slide, x, 2.82, 0.66, "", color, 12)
    text_box(slide, x - 0.38, 2.15, 1.42, 0.28, date, 11, color, True,
             align=PP_ALIGN.CENTER)
    text_box(slide, x - 0.63, 3.72, 1.92, 0.75, label, 13, INK, True,
             align=PP_ALIGN.CENTER)
rect(slide, 1.55, 5.18, 10.2, 0.92, GOLD_PALE, GOLD)
text_box(slide, 1.92, 5.45, 9.46, 0.36,
         "Later: broader ministries, website feeds, priest coordination, and additional missions",
         15, WINE_DARK, True, align=PP_ALIGN.CENTER)
notes(slide, "State clearly that the dates are targets, not promises. Security, hosting, email, and stakeholder approval are gates.")


# 10 — POC
slide = prs.slides.add_slide(blank)
header(slide, "What the proof of concept demonstrates", "A visual test of workflows and permissions—not production software.", 10)
items = [
    ("Public calendar", "Services, Confessions, Rosary, and celebrant information"),
    ("Optional signup", "Notification choices and visible ways to volunteer"),
    ("Ministry workflow", "Openings, confirmations, preferences, and substitutes"),
    ("Telegram concept", "Ministry group versus private bot interaction"),
    ("Leader review", "Balanced draft, shortages, and overrides"),
    ("Private calendar", "Father's authorized weekly view and protected fields"),
]
for i, (title, body) in enumerate(items):
    x = 0.68 + (i % 3) * 4.18
    y = 1.48 + (i // 3) * 2.25
    card(slide, x, y, 3.78, 1.78, title, body, str(i + 1), PAPER,
         GREEN if i in (0, 1) else WINE, 12)
rect(slide, 3.92, 6.08, 5.5, 0.65, WINE, WINE)
text_box(slide, 4.16, 6.25, 5.02, 0.28, "LIVE DEMONSTRATION", 13, WHITE, True,
         align=PP_ALIGN.CENTER)
notes(slide, "Transition into the live POC here. Keep the live demo focused; the PowerPoint has already established the problem and safeguards.")


# 11 — Needs
slide = prs.slides.add_slide(blank)
header(slide, "What is needed to reach alpha", "The technical work depends on a few timely decisions and inputs.", 11)
needs = [
    ("Stakeholders", "Approve a private sacristan alpha and its permission boundaries"),
    ("Webmaster", "Hosting, subdomain, authenticated email, backups, and support approach"),
    ("Pilot group", "Eight consenting sacristans, roles, preferences, and known absences"),
    ("Administrator", "Service exceptions, test sessions, and rapid feedback during August"),
]
for i, (title, body) in enumerate(needs):
    x = 0.78 + (i % 2) * 6.1
    y = 1.58 + (i // 2) * 2.22
    card(slide, x, y, 5.62, 1.72, title, body, str(i + 1), PAPER,
         WINE if i != 1 else GOLD, 13)
rect(slide, 1.45, 6.14, 10.4, 0.62, GREEN_PALE, GREEN)
text_box(slide, 1.78, 6.29, 9.74, 0.29,
         "Current scheduling remains available throughout the alpha as a fallback.",
         13, GREEN, True, align=PP_ALIGN.CENTER)
notes(slide, "This is not a request for a blank check. It is a request for a controlled test, the webmaster conversation, and access to consenting pilot users.")


# 12 — Decision
slide = prs.slides.add_slide(blank)
rect(slide, 0, 0, 13.333, 7.5, WINE_DARK, WINE_DARK, radius=False)
rect(slide, 0, 0, 13.333, 0.11, GOLD, GOLD, radius=False)
circle_badge(slide, 6.08, 0.78, 1.15, "✠", WINE, 28)
text_box(slide, 1.25, 2.20, 10.83, 0.72, "The decision requested", 32, WHITE, False,
         "Georgia", PP_ALIGN.CENTER)
text_box(slide, 1.46, 3.10, 10.42, 1.02,
         "Authorize preparation of a private sacristan alpha—subject to webmaster, privacy, and security readiness.",
         23, GOLD_PALE, True, "Georgia", PP_ALIGN.CENTER)
line(slide, 3.18, 4.48, 10.15, 4.48, GOLD, 2.2)
text_box(slide, 1.78, 4.90, 9.78, 0.55,
         "One chapel calendar. Less duplicate work. Easier participation.",
         18, WHITE, False, "Georgia", PP_ALIGN.CENTER)
text_box(slide, 4.68, 6.15, 4.0, 0.35, "QUESTIONS & DISCUSSION", 12, GOLD_PALE, True,
         align=PP_ALIGN.CENTER)
notes(slide, "Close with the specific decision, then invite discussion. If approval is conditional, record the conditions and owners before ending the meeting.")


# Basic document properties
prs.core_properties.title = "Chapel Scheduler — August 8, 2026 Planning Meeting"
prs.core_properties.subject = "High-level proposal and private alpha plan"
prs.core_properties.author = "Our Lady of Victory Chapel"
prs.core_properties.keywords = "Chapel Scheduler, scheduling, sacristans, alpha, privacy"
prs.core_properties.comments = "Prepared from the Chapel Scheduler consolidated architecture and proof of concept."

OUTPUT.parent.mkdir(parents=True, exist_ok=True)
prs.save(OUTPUT)
print(OUTPUT)
